import uuid
import requests
import os
from typing import List, Tuple, Optional, Dict
import logging
import cv2
import fitz
import shapely.geometry as sg
from shapely.geometry.base import BaseGeometry
from shapely.validation import explain_validity
from PIL import Image
from rapid_layout import RapidLayout, VisLayout
import base64
import math
from io import BytesIO
import numpy as np
import json
import subprocess
from bs4 import BeautifulSoup
from docx import Document
from xlsx2html import xlsx2html
import openpyxl
from pptx import Presentation


# ... (include all the utility functions from the previous code)

logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')

layout_engine = RapidLayout(conf_thres=0.5, model_type="pp_layout_cdla")

# pp_layout_cdla类别 (paddle OCR):
# ['text', 'title', 'figure', 'figure_caption', 'table', 'table_caption', 'header', 'footer', 'reference', 'equation']


def convert_office_to_xml(doc_bytes, type = 'slides'):
    """Convert .doc files to .docx using unoconv."""
    type_mapping = {
        "slides": "ppt",
        "document": "doc",
        "spreadsheet": "xls"
    }
    if type not in type_mapping:
        raise ValueError(f"Unsupported type: {type}")
    else :
        type_suffix = type_mapping[type]
        type_suffix_xml = f"{type_suffix}x"
    
    unique_id = str(uuid.uuid4())  # Generate a unique ID for temp files
    temp_doc_path = f"/tmp/temp_{unique_id}.{type_suffix}"
    temp_docx_path = f"/tmp/temp_{unique_id}.{type_suffix_xml}"

    # Save the .doc content to a temporary file
    with open(temp_doc_path, 'wb') as f:
        f.write(doc_bytes)

    # Convert the .doc file to .docx using unoconv
    subprocess.run(['unoconv', '-f', f'{type_suffix_xml}', temp_doc_path])

    # Read the converted .docx file
    with open(temp_docx_path, 'rb') as f:
        docx_bytes = f.read()

    # Cleanup temporary files
    os.remove(temp_doc_path)
    os.remove(temp_docx_path)

    return docx_bytes


def convert_ppt_to_pptx(doc_bytes):
    """Convert .doc files to .docx using unoconv."""
    unique_id = str(uuid.uuid4())  # Generate a unique ID for temp files
    temp_doc_path = f"/tmp/temp_{unique_id}.ppt"
    temp_docx_path = f"/tmp/temp_{unique_id}.pptx"

    # Save the .doc content to a temporary file
    with open(temp_doc_path, 'wb') as f:
        f.write(doc_bytes)

    # Convert the .doc file to .docx using unoconv
    subprocess.run(['unoconv', '-f', 'pptx', temp_doc_path])

    # Read the converted .docx file
    with open(temp_docx_path, 'rb') as f:
        docx_bytes = f.read()

    # Cleanup temporary files
    os.remove(temp_doc_path)
    os.remove(temp_docx_path)

    return docx_bytes


def pptx_to_json(pptx_bytes) -> str:
    prs = Presentation(BytesIO(pptx_bytes))
    pages = []

    # Iterate through slides
    for slide_index, slide in enumerate(prs.slides):
        page_content = {
            "page": slide_index + 1,
            "text": "",
            "table": [],
            "figures": []
        }

        # Create a BeautifulSoup object for HTML
        slide_html = BeautifulSoup("<div></div>", "html.parser")
        
        # Collect text and images
        for shape in slide.shapes:
            # Collect text, including titles and bullet points
            if hasattr(shape, "text") and shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    para_text = ' '.join(run.text for run in paragraph.runs)
                    page_content["text"] += para_text + "\n"  # Add newline for separation
                    # Create HTML paragraph for each bullet point
                    p_tag = slide_html.new_tag("p")
                    p_tag.string = para_text
                    slide_html.div.append(p_tag)
            
            # Collect images
            if shape.shape_type == 13:  # Image type
                img = shape.image
                img_bytes = img.blob
                img_base64 = base64.b64encode(img_bytes).decode('utf-8')
                img_format = img.ext
                img_data_url = f"data:image/{img_format};base64,{img_base64}"

                page_content["figures"].append(img_data_url)
                img_tag = slide_html.new_tag("img", src=img_data_url, alt="Slide Image")
                slide_html.div.append(img_tag)

        # Collect tables
        for shape in slide.shapes:
            if shape.has_table:
                table_html = slide_html.new_tag("table")
                for row in shape.table.rows:
                    tr_tag = slide_html.new_tag("tr")
                    for cell in row.cells:
                        td_tag = slide_html.new_tag("td")
                        td_tag.string = cell.text
                        tr_tag.append(td_tag)
                    table_html.append(tr_tag)
                page_content["table"].append(str(table_html))

        # Add the slide's HTML to the page content
        page_content["text"] = slide_html.div.decode()

        # Append the page data to the list
        pages.append(page_content)

    # Convert to JSON
    json_output = json.dumps(pages, indent=4)
    # Return the JSON
    # json_output_dict = json.loads(json_output)

    return json_output

def clear_style(input: str) -> str:
    # Parse the HTML input
    soup = BeautifulSoup(input, 'html.parser')

    # Remove all <style> tags
    for style in soup.find_all('style'):
        style.decompose()

    # Remove specific attributes from elements
    remove_tags = ['style', 'id', 'cellpadding', 'cellspacing', 'border']
    for tag in soup.find_all(True):
        for attr in remove_tags:
            if tag.has_attr(attr):
                del tag[attr]
        
    # Return the cleaned HTML as a string
    return str(soup)


def covert_xlsx_to_html(xlsx_bytes, clean_style:bool =True):
    unique_id = str(uuid.uuid4())
    xlsx = f"/tmp/temp_{unique_id}.xlsx"
    with open(xlsx, 'wb') as f:
        f.write(xlsx_bytes)
    xlsx_ostream = xlsx2html(xlsx)
    xlsx_ostream.seek(0)
    context = ''
    for line in xlsx_ostream:
        context += line
    if clean_style:
        context = clear_style(context)

    return context

def covert_xlsx_to_html2(xlsx_bytes, clean_style:bool =True):
    unique_id = str(uuid.uuid4())
    xlsx = f"/tmp/temp_{unique_id}.xlsx"
    with open(xlsx, 'wb') as f:
        f.write(xlsx_bytes)

    sheet_names = openpyxl.load_workbook(xlsx).sheetnames
    
    sheet_number = len(sheet_names)
    context = ''

    for sheet_number_inx in range(sheet_number):
        xlsx_ostream = xlsx2html(xlsx, sheet = sheet_number_inx)
        xlsx_ostream.seek(0)
        context = context + f'<h1>{sheet_names[sheet_number_inx]}</h1>'
        for line in xlsx_ostream:
            context += line

    if clean_style:
        context = clear_style(context)
    return context


def convert_docx_to_markdown(docx_bytes):
    """Convert .docx files to markdown."""
    document = Document(BytesIO(docx_bytes))
    markdown = ""
    for para in document.paragraphs:
        markdown += para.text + "\n\n"
    return markdown.strip()


def parser_lv1(input: str, type = "document") -> str:
    """Parse the input base64 string and return markdown in base64."""
    doc_bytes = base64.b64decode(input)

    if type == "document":
        # Check file type based on magic numbers
        if doc_bytes[:4] == b'\xD0\xCF\x11\xE0':  # DOC magic number
            docx_bytes = convert_office_to_xml(doc_bytes, type)
            markdown = convert_docx_to_markdown(docx_bytes)
        
        if doc_bytes[:4] == b'\x50\x4B\x03\x04': # docx magic number
            markdown = convert_docx_to_markdown(doc_bytes)
    
    if type == "spreadsheet":
        
        if doc_bytes[:4] == b'\xD0\xCF\x11\xE0': # xls magic number
            docx_bytes = convert_office_to_xml(doc_bytes, type)
            html = covert_xlsx_to_html(docx_bytes)
            markdown = html
        
        if doc_bytes[:4] == b'\x50\x4B\x03\x04': # xlsx magic number
            html = covert_xlsx_to_html(doc_bytes)
            markdown = html
    
    if type == 'slides':

        if doc_bytes[:4] == b'\xD0\xCF\x11\xE0':
            doc_bytes = convert_office_to_xml(doc_bytes, type)
            markdown = pptx_to_json(doc_bytes)

        if doc_bytes[:4] == b'\x50\x4B\x03\x04':
            markdown = pptx_to_json(doc_bytes)



    context_base64 = base64.b64encode(markdown.encode()).decode()
    result = {
        "context_base64": context_base64
    }
    return json.dumps(result)

def round_by_factor(number: int, factor: int) -> int:
    """Returns the closest integer to 'number' that is divisible by 'factor'."""
    return round(number / factor) * factor


def ceil_by_factor(number: int, factor: int) -> int:
    """Returns the smallest integer greater than or equal to 'number' that is divisible by 'factor'."""
    return math.ceil(number / factor) * factor


def floor_by_factor(number: int, factor: int) -> int:
    """Returns the largest integer less than or equal to 'number' that is divisible by 'factor'."""
    return math.floor(number / factor) * factor

def smart_resize(
    height: int, width: int, factor: int, min_pixels: int, max_pixels: int, MAX_RATIO) -> tuple[int, int]:
    """
    Rescales the image so that the following conditions are met:

    1. Both dimensions (height and width) are divisible by 'factor'.

    2. The total number of pixels is within the range ['min_pixels', 'max_pixels'].

    3. The aspect ratio of the image is maintained as closely as possible.
    """
    if max(height, width) / min(height, width) > MAX_RATIO:
        raise ValueError(
            f"absolute aspect ratio must be smaller than {MAX_RATIO}, got {max(height, width) / min(height, width)}"
        )
    h_bar = max(factor, round_by_factor(height, factor))
    w_bar = max(factor, round_by_factor(width, factor))
    if h_bar * w_bar > max_pixels:
        beta = math.sqrt((height * width) / max_pixels)
        h_bar = floor_by_factor(height / beta, factor)
        w_bar = floor_by_factor(width / beta, factor)
    elif h_bar * w_bar < min_pixels:
        beta = math.sqrt(min_pixels / (height * width))
        h_bar = ceil_by_factor(height * beta, factor)
        w_bar = ceil_by_factor(width * beta, factor)
    return h_bar, w_bar

def extract_vision_info(conversations: list[dict] | list[list[dict]]) -> list[dict]:
    vision_infos = []
    if isinstance(conversations[0], dict):
        conversations = [conversations]
    for conversation in conversations:
        for message in conversation:
            if isinstance(message["content"], list):
                for ele in message["content"]:
                    if (
                        "image" in ele
                        or "image_url" in ele
                        or "video" in ele
                        or ele["type"] in ("image", "image_url", "video")
                    ):
                        vision_infos.append(ele)
    return vision_infos

def fetch_image(ele: dict[str, str | Image.Image], IMAGE_FACTOR, MIN_PIXELS, MAX_PIXELS, MAX_RATIO) -> Image.Image:
    size_factor = IMAGE_FACTOR
    if "image" in ele:
        image = ele["image"]
    else:
        image = ele["image_url"]
    image_obj = None
    if isinstance(image, Image.Image):
        image_obj = image
    elif image.startswith("http://") or image.startswith("https://"):
        image_obj = Image.open(requests.get(image, stream=True).raw)
    elif image.startswith("file://"):
        image_obj = Image.open(image[7:])
    elif image.startswith("data:image"):
        data = image.split(";", 1)[1]
        if data.startswith("base64,"):
            data = base64.b64decode(data[7:])
            image_obj = Image.open(BytesIO(data))
    else:
        image_obj = Image.open(image)
    if image_obj is None:
        raise ValueError(f"Unrecognized image input, support local path, http url, base64 and PIL.Image, got {image}")
    image = image_obj.convert("RGB")
    ## resize
    if "resized_height" in ele and "resized_width" in ele:
        resized_height, resized_width = smart_resize(
            ele["resized_height"],
            ele["resized_width"],
            factor = IMAGE_FACTOR,
            min_pixels = MIN_PIXELS, 
            max_pixels = MAX_PIXELS,
            MAX_RATIO = MAX_RATIO
        )
    else:
        width, height = image.size
        min_pixels = ele.get("min_pixels", MIN_PIXELS)
        max_pixels = ele.get("max_pixels", MAX_PIXELS)
        resized_height, resized_width = smart_resize(
            height,
            width,
            factor=size_factor,
            min_pixels=min_pixels,
            max_pixels=max_pixels,
            MAX_RATIO = MAX_RATIO
        )
    image = image.resize((resized_width, resized_height))

    return image

def process_vision_info(
    conversations: list[dict] | list[list[dict]], IMAGE_FACTOR, MIN_PIXELS, MAX_PIXELS, MAX_RATIO) :
    vision_infos = extract_vision_info(conversations)
    ## Read images or videos
    image_inputs = []
    video_inputs = []
    for vision_info in vision_infos:
        if "image" in vision_info or "image_url" in vision_info:
            image_inputs.append(fetch_image(vision_info, IMAGE_FACTOR, MIN_PIXELS, MAX_PIXELS, MAX_RATIO))
        else:
            raise ValueError("image, image_url should in content.")
    if len(image_inputs) == 0:
        image_inputs = None
    return image_inputs, video_inputs


def callvlm(img_data, vlm_api_url, vlm_api_key, system_prompts):
    
    url = vlm_api_url
    headers = {"Content-Type": "application/json", "Authorization": f"Bearer {vlm_api_key}"}
    data = { 
        "model": "atom",  
        "stream": False, 
        "messages": [ {
            "role": "user",
            "content": [ {
                "text": f'''{system_prompts}''',
                 "type": "text"}, 
                 { "type": "image_url", 
                  "image_url": {
                      "url": f'''data:image/jpeg;base64,{img_data}'''}}]
                  }]}
    response = requests.post(url=url, headers=headers, data=json.dumps(data))
    return response.text

def parse_pdf_to_images(pdf_path: str, 
                         output_dir: str = './output') -> List[Tuple[str, List[str]]]:
    image_infos = []
    pdf_document = fitz.open(pdf_path)
    for page_index, page in enumerate(pdf_document):
        rect_images = []
        logging.info(f'parse page: {page_index}')
        # 保存页面为图片
        pix = page.get_pixmap(matrix=fitz.Matrix(4, 4))
        pix = Image.frombytes('RGB', [pix.width, pix.height], pix.samples)
        boxes, scores, class_names, elapse = layout_engine(pix)
        for index, (class_name, box) in enumerate(zip(class_names, boxes)):
            if class_name == 'figure' or class_name == 'table':
                name = f'{page_index}_{index}.png'
                sub_pix = pix.crop(box)
                sub_pix.save(os.path.join(output_dir, name))
                rect_images.append(name)

        boxes_ = []
        scores_ = []
        class_names_ = []
        for i, (class_name, box, score) in enumerate(zip(class_names, boxes, scores)):
            if class_name == 'figure' or class_name == 'table':
                boxes_.append(box)
                scores_.append(score)
                class_name = f'{page_index}_{i}.png'
                class_names_.append(class_name)
                
        page_image = os.path.join(output_dir, f'{page_index}.png')
        pix = np.array(pix)
        pix = cv2.cvtColor(pix, cv2.COLOR_RGB2BGR)
        print(boxes_, scores_, class_names_)
        ploted_img = VisLayout.draw_detections(pix, boxes_, scores_, class_names_)
        if ploted_img is not None:
            cv2.imwrite(page_image, ploted_img)
        # ploted_img.save(page_image)
        image_infos.append((page_image, rect_images))
    pdf_document.close()
    return image_infos

def image_to_base64(image: Image.Image) -> str:
    # Create a BytesIO object to hold the image data
    buffered = BytesIO()
    
    # Save the image to the BytesIO object in PNG format
    image.save(buffered, format="JPEG")
    
    # Get the byte data from the BytesIO object
    img_bytes = buffered.getvalue()
    
    # Encode the byte data to Base64
    img_base64 = base64.b64encode(img_bytes).decode('utf-8')
    
    return img_base64

def _gpt_parse_images(
        image_infos: List[Tuple[str, List[str]]],
        DEFAULT_PROMPT,
        DEFAULT_RECT_PROMPT,
        DEFAULT_ROLE_PROMPT,
        IMAGE_FACTOR, MIN_PIXELS, MAX_PIXELS, MAX_RATIO,
        vlm_api_url, vlm_api_key, system_prompts,
        prompt_dict: Optional[Dict] = None,
        output_dir: str = './',
        api_key: Optional[str] = None,
        base_url: Optional[str] = None,
        model: str = 'gpt-4o',
        verbose: bool = False,
        gpt_worker: int = 1,

):
    """
    Parse images to markdown content.
    """

    if isinstance(prompt_dict, dict) and 'prompt' in prompt_dict:
        prompt = prompt_dict['prompt']
        logging.info("prompt is provided, using user prompt.")
    else:
        prompt = DEFAULT_PROMPT
        logging.info("prompt is not provided, using default prompt.")
    if isinstance(prompt_dict, dict) and 'rect_prompt' in prompt_dict:
        rect_prompt = prompt_dict['rect_prompt']
        logging.info("rect_prompt is provided, using user prompt.")
    else:
        rect_prompt = DEFAULT_RECT_PROMPT
        logging.info("rect_prompt is not provided, using default prompt.")
    if isinstance(prompt_dict, dict) and 'role_prompt' in prompt_dict:
        role_prompt = prompt_dict['role_prompt']
        logging.info("role_prompt is provided, using user prompt.")
    else:
        role_prompt = DEFAULT_ROLE_PROMPT
        logging.info("role_prompt is not provided, using default prompt.")

    def _process_page(index: int, image_info: Tuple[str, List[str]], **args) -> Tuple[int, str]:
        logging.info(f'gpt parse page: {index}')

        # agent = Agent(role=role_prompt, api_key=api_key, base_url=base_url, disable_python_run=True, model=model, **args)
        page_image, rect_images = image_info
        local_prompt = prompt
        local_prompt = role_prompt + local_prompt
        if rect_images:
            local_prompt += rect_prompt + ', '.join(rect_images)
        
        messages = [{
        "role": "user",
        "content": [
            {
                "type": "image",
                "image": page_image,
            },
            {"type": "text", "text": local_prompt},
        ],}]
        image_inputs, video_inputs = process_vision_info(messages, IMAGE_FACTOR, MIN_PIXELS, MAX_PIXELS, MAX_RATIO)
        
        output_text = callvlm(image_to_base64(image_inputs[0]), vlm_api_url, vlm_api_key, system_prompts)

        return index, output_text

    contents = [None] * len(image_infos)

    for index, singl_img_page_path in enumerate(image_infos):
        indexer, output_text = _process_page(index, singl_img_page_path)
        contents[indexer] = output_text

    return contents


